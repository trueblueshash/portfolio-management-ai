"""
Document processing service: extract text, chunk, embed, summarize.
"""
import logging
import os
from typing import Optional
from sqlalchemy.orm import Session
from app.models.document import PortfolioDocument, DocumentChunk
from openai import OpenAI
from app.core.config import settings
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)

# OpenAI client for embeddings
# Note: For embeddings, we can use OpenAI API or OpenRouter
# Using OpenRouter with OpenAI-compatible format
openai_client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=settings.OPENROUTER_API_KEY,
)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page.extract_text()
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        doc = DocxDocument(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise


def extract_text(file_path: str, mime_type: str, google_doc_id: Optional[str] = None) -> str:
    """Extract text from document based on MIME type"""
    if mime_type == "application/vnd.google-apps.document":
        # Handle Google Docs - lazy import to avoid circular dependency
        from app.services.gdocs_service import extract_google_doc_content, extract_google_doc_id
        
        if google_doc_id:
            doc_id = google_doc_id
        else:
            # Extract doc ID from file_path if it's a URL
            if file_path.startswith("http"):
                doc_id = extract_google_doc_id(file_path)
            else:
                # Assume file_path is the doc ID itself
                doc_id = file_path
        
        doc_content = extract_google_doc_content(doc_id)
        return doc_content.get('full_text', '')
    elif mime_type == "application/pdf":
        return extract_text_from_pdf(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported MIME type: {mime_type}")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[tuple[str, int]]:
    """
    Chunk text into overlapping segments.
    Returns list of (chunk_text, page_number) tuples.
    """
    chunks = []
    
    # Check if text has page/slide markers
    has_page_markers = "--- Page" in text or "--- Slide" in text
    
    if has_page_markers:
        # Has page markers - split and process by page
        if "--- Page" in text:
            pages = text.split("--- Page")
        else:
            pages = text.split("--- Slide")
        
        current_page = 1
        for page_content in pages[1:]:  # Skip first empty split
            # Extract page number
            if "---" in page_content:
                page_num_str = page_content.split("---")[0].strip()
                try:
                    current_page = int(page_num_str)
                except:
                    pass
                page_text = "---".join(page_content.split("---")[1:])
            else:
                page_text = page_content
            
            # Chunk this page
            page_chunks = _chunk_string(page_text, chunk_size, overlap)
            for chunk in page_chunks:
                chunks.append((chunk, current_page))
    else:
        # No page markers - chunk entire text normally
        text_chunks = _chunk_string(text, chunk_size, overlap)
        for chunk in text_chunks:
            chunks.append((chunk, None))
    
    return chunks


def _chunk_string(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Helper to chunk a string"""
    if not text or not text.strip():
        return []
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end].strip()
        
        # Only add non-empty chunks
        if chunk:
            chunks.append(chunk)
        
        if end >= text_len:
            break
        
        # Move start forward by chunk_size - overlap
        start += chunk_size - overlap
    
    # If no chunks were created (shouldn't happen, but safety check)
    if not chunks and text.strip():
        chunks.append(text.strip()[:chunk_size])
    
    return chunks


def chunk_google_doc_by_structure(doc_content: dict, sections: list[dict]) -> list[tuple[str, Optional[str], dict]]:
    """
    Smart chunking for structured Google Docs.
    
    Instead of arbitrary 800-char chunks:
    1. Chunk by section (Nov'25, Dec'25, etc.)
    2. Keep subsections together (Summary, Metrics, Q&A)
    3. Preserve context within quarters
    4. Add metadata: {section: "Nov'25", subsection: "Summary"}
    
    Returns:
    List of (chunk_text, source_section, metadata) tuples
    """
    chunks = []
    
    for section in sections:
        section_heading = section.get("heading", "")
        subsections = section.get("subsections", {})
        section_content = section.get("content", "")
        
        # Chunk each subsection separately
        for subsection_name, subsection_content in subsections.items():
            if not subsection_content or not subsection_content.strip():
                continue
            
            # Create chunk for this subsection
            chunk_text = f"{section_heading} - {subsection_name}\n\n{subsection_content.strip()}"
            metadata = {
                "section": section_heading,
                "subsection": subsection_name,
                "contains_metrics": subsection_name.lower() == "metrics",
                "has_qa": subsection_name.lower() in ["q&a", "qa", "questions to ask"]
            }
            chunks.append((chunk_text, section_heading, metadata))
        
        # Also create a chunk for the full section content if it exists and has content
        if section_content and section_content.strip():
            chunk_text = f"{section_heading}\n\n{section_content.strip()}"
            metadata = {
                "section": section_heading,
                "subsection": "full",
                "contains_metrics": "metrics" in section_content.lower(),
                "has_qa": "q&a" in section_content.lower()
            }
            chunks.append((chunk_text, section_heading, metadata))
        
        # If no subsections and no content, but we have a heading, create a minimal chunk
        if not subsections and not section_content and section_heading:
            chunk_text = section_heading
            metadata = {
                "section": section_heading,
                "subsection": "heading_only",
                "contains_metrics": False,
                "has_qa": False
            }
            chunks.append((chunk_text, section_heading, metadata))
    
    return chunks


def generate_embedding(text: str) -> list[float]:
    """Generate embedding using OpenAI via OpenRouter"""
    try:
        # Use OpenAI's embedding model via OpenRouter
        response = openai_client.embeddings.create(
            model="openai/text-embedding-3-small",
            input=text[:8000]  # Limit text length for embeddings
        )
        return response.data[0].embedding
    except Exception as e:
        logger.error(f"Error generating embedding: {e}")
        # Fallback: return empty embedding
        return [0.0] * 1536


def generate_summary(text: str, company_name: str = None) -> str:
    """Generate document summary using Claude Haiku"""
    from openai import OpenAI
    from app.core.config import settings
    
    client = OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=settings.OPENROUTER_API_KEY,
    )
    
    # Truncate text for summary (first 5000 chars)
    text_preview = text[:5000]
    
    prompt = f"""Create a concise 2-3 paragraph summary of this document.

{"Company: " + company_name if company_name else ""}

Focus on:
- Key topics and themes
- Important decisions or recommendations
- Notable metrics or data points
- Main conclusions

Document content:
{text_preview}

Summary:"""

    try:
        response = client.chat.completions.create(
            model="anthropic/claude-3.5-haiku",
            messages=[
                {
                    "role": "system",
                    "content": "You are a document summarizer. Create clear, informative summaries."
                },
                {"role": "user", "content": prompt}
            ],
            max_tokens=300,
            temperature=0.3,
        )
        
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error generating summary: {e}")
        return "Summary generation failed."


def process_document(document_id: str, db: Session):
    """
    Process a document: extract text, chunk, embed, summarize.
    Returns True if successful, False otherwise.
    """
    from uuid import UUID
    
    try:
        doc_uuid = UUID(document_id)
    except ValueError:
        raise ValueError(f"Invalid document_id format: {document_id}")
    
    document = db.query(PortfolioDocument).filter(PortfolioDocument.id == doc_uuid).first()
    if not document:
        raise ValueError(f"Document not found: {document_id}")
    
    logger.info(f"Processing document: {document.title}")
    
    try:
        # 1. Extract text
        logger.info(f"  Extracting text...")
        if document.mime_type == 'application/vnd.google-apps.document':
            if not document.google_doc_id:
                raise ValueError(f"Google Doc document {document.id} missing google_doc_id")
            logger.info(f"  Extracting from Google Doc: {document.google_doc_id}")
            full_text = extract_text(document.file_path or "", document.mime_type, google_doc_id=document.google_doc_id)
        else:
            if not document.file_path:
                raise ValueError(f"Document {document.id} missing file_path")
            logger.info(f"  Extracting from file: {document.file_path}")
            full_text = extract_text(document.file_path, document.mime_type)
        
        # Validate extracted text
        if not full_text or len(full_text.strip()) == 0:
            logger.error(f"  ❌ No text extracted from document")
            document.is_processed = False
            db.commit()
            return False
        
        logger.info(f"  Extracted {len(full_text)} characters")
        document.full_text = full_text
        
        # Delete existing chunks before creating new ones
        logger.info(f"  Deleting existing chunks...")
        db.query(DocumentChunk).filter(DocumentChunk.document_id == document.id).delete()
        db.flush()  # Flush to ensure delete is applied
        
        # 2. Chunk text
        logger.info(f"  Chunking text...")
        company = document.company
        
        # Use structured chunking for Google Docs that are primary sources
        if document.mime_type == 'application/vnd.google-apps.document' and document.is_primary_source:
            # Get structured content for intelligent chunking - lazy import to avoid circular dependency
            from app.services.gdocs_service import extract_google_doc_content
            
            if not document.google_doc_id:
                raise ValueError(f"Google Doc document {document.id} missing google_doc_id")
            doc_content = extract_google_doc_content(document.google_doc_id)
            sections = doc_content.get('sections', [])
            
            # Only use structured chunking if sections exist
            if sections and len(sections) > 0:
                structured_chunks = chunk_google_doc_by_structure(doc_content, sections)
                # structured_chunks returns: (chunk_text, source_section, metadata)
                # Convert to format compatible with chunk_text: (chunk_text, page_num)
                chunks = [(chunk_text_val, None) for chunk_text_val, source_section, _ in structured_chunks]
                logger.info(f"  Used structured chunking: {len(chunks)} chunks")
            else:
                # No sections found, use regular chunking
                chunks = chunk_text(full_text)
                logger.info(f"  No sections found, used regular chunking: {len(chunks)} chunks")
        else:
            # Regular chunking - CORRECT CALL
            # chunk_text returns: list[tuple[str, int]] = [(chunk_text, page_num), ...]
            chunks = chunk_text(full_text)
            logger.info(f"  Used regular chunking: {len(chunks)} chunks")
        
        # Validate chunks
        if len(chunks) == 0:
            logger.error(f"  ❌ Chunking returned 0 chunks (text length: {len(full_text)})")
            document.is_processed = False
            db.commit()
            return False
        
        logger.info(f"  Created {len(chunks)} chunks")
        
        # 3. Generate embeddings and save chunks
        logger.info(f"  Generating embeddings...")
        chunks_created = 0
        for i, (chunk_text_str, page_num) in enumerate(chunks):
            try:
                if not chunk_text_str or not chunk_text_str.strip():
                    logger.warning(f"  Skipping empty chunk {i}")
                    continue
                
                embedding = generate_embedding(chunk_text_str)
                
                # Use doc_type as source_section for non-Google Docs
                source_section = document.doc_type if document.doc_type else "unknown"
                
                chunk_obj = DocumentChunk(
                    document_id=document.id,
                    chunk_index=i,
                    chunk_text=chunk_text_str,
                    chunk_embedding=embedding,  # Store as vector array (pgvector handles conversion)
                    page_number=page_num,
                    source_section=source_section
                )
                db.add(chunk_obj)
                chunks_created += 1
            except Exception as e:
                logger.error(f"  Error processing chunk {i}: {e}", exc_info=True)
                continue
        
        # Commit chunks before continuing
        db.commit()
        logger.info(f"  Saved {chunks_created} chunks to database")
        
        if chunks_created == 0:
            logger.error(f"  ❌ No chunks were saved")
            document.is_processed = False
            db.commit()
            return False
        
        # 4. Generate summary
        logger.info(f"  Generating summary...")
        company_name = company.name if company else None
        summary = generate_summary(full_text, company_name)
        document.summary = summary
        
        # 5. Mark as processed
        document.is_processed = True
        db.commit()
        
        logger.info(f"  ✅ Document processed successfully: {chunks_created} chunks created")
        return True
        
    except Exception as e:
        logger.error(f"  ❌ Error processing document: {e}", exc_info=True)
        document.is_processed = False
        db.rollback()
        db.commit()
        raise


def process_document_background(document_id: str):
    """
    Wrapper for process_document that creates its own database session.
    Used with FastAPI BackgroundTasks.
    """
    from app.db.session import SessionLocal
    
    db = SessionLocal()
    try:
        process_document(document_id, db)
    finally:
        db.close()
